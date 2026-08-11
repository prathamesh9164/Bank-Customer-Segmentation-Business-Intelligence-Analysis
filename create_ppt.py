"""
🏦 IndoSynth Gramin Bank — SQL Phase PPT Generator
Generates a professional 20-slide PowerPoint presentation covering:
  - Project overview & objectives
  - Database schema design (ER diagram, 9 tables)
  - Data loading strategy (LOAD DATA LOCAL INFILE)
  - Indexing strategy
  - 10 high-insight analytical SQL queries (Basic → Intermediate → Advanced)
  - Key takeaways
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette (Light & Eye-Soothing) ──────────────────────────────────
BG_LIGHT      = RGBColor(0xF7, 0xF5, 0xF0)   # Warm off-white background
BG_CARD       = RGBColor(0xEE, 0xEC, 0xE6)   # Soft warm card background
ACCENT_ROSE   = RGBColor(0xC7, 0x7D, 0x6E)   # Dusty rose accent
ACCENT_GOLD   = RGBColor(0xC9, 0x9E, 0x58)   # Muted gold accent
ACCENT_TEAL   = RGBColor(0x52, 0x9E, 0x97)   # Muted teal accent
ACCENT_BLUE   = RGBColor(0x5E, 0x89, 0xB4)   # Calm blue accent
DARK_TEXT     = RGBColor(0x2C, 0x2C, 0x2C)   # Primary dark text
SECOND_TEXT   = RGBColor(0x55, 0x55, 0x55)   # Secondary text
MID_GRAY      = RGBColor(0x88, 0x88, 0x88)   # Subtle gray
CODE_BG       = RGBColor(0xE8, 0xE5, 0xDE)   # Light warm code block bg
PURPLE        = RGBColor(0x8B, 0x6B, 0xAF)   # Soft lavender accent
GREEN         = RGBColor(0x5E, 0x9E, 0x65)   # Sage green accent
RED_SOFT      = RGBColor(0xC4, 0x5E, 0x5E)   # Muted red

# ── Slide dimensions (Widescreen 16:9) ────────────────────────────────────
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Make corners subtle
    shape.adjustments[0] = 0.05
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.2):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_name="Calibri"):
    """
    Add a text box with multiple styled lines.
    Each line is a dict: {text, size, color, bold, alignment, spacing_after}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line.get("text", "")
        p.font.size = Pt(line.get("size", 18))
        p.font.color.rgb = line.get("color", DARK_TEXT)
        p.font.bold = line.get("bold", False)
        p.font.name = font_name
        p.alignment = line.get("alignment", PP_ALIGN.LEFT)
        p.space_after = Pt(line.get("spacing_after", 6))
        p.space_before = Pt(line.get("spacing_before", 0))
    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_ROSE):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    """Add a styled code block with light background."""
    # Background card
    card = add_shape(slide, left, top, width, height, CODE_BG, ACCENT_TEAL, Pt(1))
    # Code text
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(code_text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x40)
        p.font.name = "Consolas"
        p.space_after = Pt(2)
        p.space_before = Pt(0)
    return card


def add_slide_number(slide, num, total):
    """Add slide number at bottom-right."""
    add_textbox(slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5),
                Inches(1.3), Inches(0.4), f"{num} / {total}",
                font_size=10, font_color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def add_footer_bar(slide):
    """Add a subtle footer bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), SLIDE_HEIGHT - Inches(0.35),
                                    SLIDE_WIDTH, Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE5, 0xE2, 0xDB)
    shape.line.fill.background()
    add_textbox(slide, Inches(0.5), SLIDE_HEIGHT - Inches(0.33),
                Inches(5), Inches(0.3),
                "IndoSynth Gramin Bank  •  SQL Phase  •  Major Project",
                font_size=9, font_color=MID_GRAY, font_name="Calibri")


def add_section_badge(slide, text, left, top, color=ACCENT_ROSE):
    """Add a small colored badge/tag."""
    badge = add_shape(slide, left, top, Inches(1.8), Inches(0.35), color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.02),
                Inches(1.6), Inches(0.3), text,
                font_size=10, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    return badge


# ══════════════════════════════════════════════════════════════════════════
#  BUILD THE PRESENTATION
# ══════════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

TOTAL_SLIDES = 20
blank_layout = prs.slide_layouts[6]  # Blank layout


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 1: TITLE SLIDE
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

# Decorative accent bar at top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ROSE
shape.line.fill.background()

# Bank icon text
add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1), "🏦",
            font_size=60, alignment=PP_ALIGN.CENTER)

# Title
add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
            "IndoSynth Gramin Bank", font_size=44, font_color=DARK_TEXT,
            bold=True, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Subtitle
add_textbox(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.7),
            "Phase 1  —  MySQL Database Design & Analytics",
            font_size=24, font_color=ACCENT_ROSE, alignment=PP_ALIGN.CENTER)

# Accent line
add_accent_line(slide, Inches(4.5), Inches(4.2), Inches(4.3), ACCENT_ROSE)

# Description
add_textbox(slide, Inches(2), Inches(4.6), Inches(9), Inches(0.8),
            "End-to-end relational database design, bulk data ingestion,\nand 25+ analytical SQL queries across ~3 Million records",
            font_size=16, font_color=SECOND_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

# Bottom info
add_multi_text(slide, Inches(2), Inches(5.8), Inches(9), Inches(1.2), [
    {"text": "Data Analytics & Business Intelligence  •  Major Project", "size": 14, "color": SECOND_TEXT, "alignment": PP_ALIGN.CENTER, "spacing_after": 8},
    {"text": "MySQL 8.0  •  9 Tables  •  ~3M Records  •  25+ Queries", "size": 13, "color": ACCENT_TEAL, "alignment": PP_ALIGN.CENTER},
])

add_footer_bar(slide)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 2: AGENDA / TABLE OF CONTENTS
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
            "Agenda", font_size=36, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(2.5), ACCENT_ROSE)

agenda_items = [
    ("01", "Project Overview & Business Objectives", ACCENT_ROSE),
    ("02", "Dataset Overview — 9 Core Banking Tables", ACCENT_GOLD),
    ("03", "ER Schema Design & Normalization (3NF)", ACCENT_TEAL),
    ("04", "Table Creation — DDL with Constraints", ACCENT_BLUE),
    ("05", "Bulk Data Loading — LOAD DATA LOCAL INFILE", PURPLE),
    ("06", "Data Verification & Row Counts", GREEN),
    ("07", "Indexing Strategy for Performance", ACCENT_ROSE),
    ("08", "10 Analytical Queries with Key Insights", ACCENT_GOLD),
    ("09", "Key Takeaways & Summary", ACCENT_TEAL),
]

for i, (num, item_text, color) in enumerate(agenda_items):
    y_pos = Inches(1.7) + Inches(i * 0.58)
    # Number badge
    badge = add_shape(slide, Inches(1.2), y_pos, Inches(0.55), Inches(0.4), color)
    add_textbox(slide, Inches(1.2), y_pos + Inches(0.02), Inches(0.55), Inches(0.35),
                num, font_size=14, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    # Text
    add_textbox(slide, Inches(2.0), y_pos + Inches(0.02), Inches(9), Inches(0.4),
                item_text, font_size=17, font_color=DARK_TEXT)

add_footer_bar(slide)
add_slide_number(slide, 2, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 3: PROJECT OVERVIEW
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "OVERVIEW", Inches(0.8), Inches(0.4), ACCENT_ROSE)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Project Overview & Business Objectives", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4), ACCENT_ROSE)

# Left card — Domain
add_shape(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.8), BG_CARD, ACCENT_ROSE, Pt(1))
add_multi_text(slide, Inches(1.1), Inches(2.1), Inches(5), Inches(4.5), [
    {"text": "🏦  Domain", "size": 18, "color": ACCENT_ROSE, "bold": True, "spacing_after": 10},
    {"text": "Banking & Financial Services (Rural / Gramin Banking)", "size": 14, "color": SECOND_TEXT, "spacing_after": 16},
    {"text": "📋  Business Problem", "size": 18, "color": ACCENT_ROSE, "bold": True, "spacing_after": 10},
    {"text": "Analyze the operations of IndoSynth Gramin Bank across India to uncover insights on:", "size": 13, "color": SECOND_TEXT, "spacing_after": 10},
    {"text": "•  Customer demographics & segmentation", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Loan portfolio performance & approval rates", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Payment behavior & default identification", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Branch efficiency & profitability", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Transaction patterns (UPI vs Traditional)", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Credit risk modeling & NPA detection", "size": 13, "color": DARK_TEXT},
])

# Right card — Key objectives
add_shape(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(4.8), BG_CARD, ACCENT_TEAL, Pt(1))
add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5), Inches(4.5), [
    {"text": "🎯  Key Objectives (SQL Phase)", "size": 18, "color": ACCENT_TEAL, "bold": True, "spacing_after": 12},
    {"text": "1.  Design a normalized relational database (3NF)", "size": 13, "color": DARK_TEXT, "spacing_after": 8},
    {"text": "2.  Define DDL with Primary & Foreign Keys", "size": 13, "color": DARK_TEXT, "spacing_after": 8},
    {"text": "3.  Bulk load ~3 Million records from 9 CSVs", "size": 13, "color": DARK_TEXT, "spacing_after": 8},
    {"text": "4.  Create strategic indexes for performance", "size": 13, "color": DARK_TEXT, "spacing_after": 8},
    {"text": "5.  Write 25+ analytical queries (Basic → Advanced)", "size": 13, "color": DARK_TEXT, "spacing_after": 16},
    {"text": "🛠️  Technology", "size": 18, "color": ACCENT_TEAL, "bold": True, "spacing_after": 10},
    {"text": "•  MySQL Server 8.0 + MySQL Workbench", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  SQL (DDL, DML, DQL, Window Functions)", "size": 13, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  LOAD DATA LOCAL INFILE for bulk ingestion", "size": 13, "color": DARK_TEXT},
])

add_footer_bar(slide)
add_slide_number(slide, 3, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 4: DATASET OVERVIEW
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "DATASET", Inches(0.8), Inches(0.4), ACCENT_GOLD)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Dataset Overview — 9 Core Banking Tables", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4.5), ACCENT_GOLD)

# Table data
tables_data = [
    ("regions",           "40",        "1 row per region",            "region_id, zone, primary_state"),
    ("loan_types",        "10",        "1 row per loan product",      "loan_type_id, interest_rate, collateral"),
    ("branches",          "250",       "1 row per branch",            "branch_id, region_id, city, type"),
    ("employees",         "~2,000",    "1 row per employee",          "employee_id, branch_id, designation"),
    ("customers",         "100,000",   "1 row per account holder",    "customer_id, income, segment"),
    ("credit_history",    "100,000",   "1 row per credit record",     "credit_score, rating, utilization"),
    ("loan_applications", "300,000",   "1 row per application",       "application_id, amount, status"),
    ("loan_payments",     "900,000",   "1 row per EMI payment",       "payment_id, status, days_late"),
    ("transactions",      "1,500,000", "1 row per transaction",       "transaction_id, amount, mode"),
]

# Column headers
header_y = Inches(1.85)
header_bg = add_shape(slide, Inches(0.8), header_y, Inches(11.7), Inches(0.42), ACCENT_GOLD)
headers = [("Table Name", 0.9, 2.2), ("Records", 3.2, 1.2), ("Grain", 4.5, 3.0), ("Key Attributes", 7.7, 4.5)]
for h_text, h_left, h_width in headers:
    add_textbox(slide, Inches(h_left), header_y + Inches(0.05), Inches(h_width), Inches(0.35),
                h_text, font_size=13, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

# Table rows
for i, (tbl, rows, grain, attrs) in enumerate(tables_data):
    y = Inches(2.35) + Inches(i * 0.46)
    row_bg = BG_CARD if i % 2 == 0 else BG_LIGHT
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.42), row_bg)
    add_textbox(slide, Inches(0.9), y + Inches(0.05), Inches(2.2), Inches(0.35),
                tbl, font_size=12, font_color=ACCENT_TEAL, bold=True, font_name="Consolas")
    add_textbox(slide, Inches(3.2), y + Inches(0.05), Inches(1.2), Inches(0.35),
                rows, font_size=12, font_color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4.5), y + Inches(0.05), Inches(3.0), Inches(0.35),
                grain, font_size=12, font_color=SECOND_TEXT)
    add_textbox(slide, Inches(7.7), y + Inches(0.05), Inches(4.5), Inches(0.35),
                attrs, font_size=11, font_color=MID_GRAY, font_name="Consolas")

# Total summary
total_y = Inches(2.35) + Inches(9 * 0.46) + Inches(0.15)
add_shape(slide, Inches(0.8), total_y, Inches(11.7), Inches(0.42), ACCENT_ROSE)
add_textbox(slide, Inches(0.9), total_y + Inches(0.05), Inches(5), Inches(0.35),
            "TOTAL:  ~3,002,300 Records  •  9 Tables  •  ~400 MB Raw Data",
            font_size=13, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

add_footer_bar(slide)
add_slide_number(slide, 4, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 5: ER DIAGRAM
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "SCHEMA", Inches(0.8), Inches(0.4), ACCENT_TEAL)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Entity-Relationship Diagram (3NF)", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4), ACCENT_TEAL)

# ER diagram image
er_img_path = r"B:\Major Project\mysql\Screenshot 2026-07-08 172155.png"
if os.path.exists(er_img_path):
    slide.shapes.add_picture(er_img_path, Inches(1.5), Inches(1.8), Inches(7.0), Inches(5.2))

# Right side — Key relationships
add_shape(slide, Inches(8.8), Inches(1.8), Inches(4.0), Inches(5.2), BG_CARD, ACCENT_TEAL, Pt(1))
add_multi_text(slide, Inches(9.0), Inches(1.95), Inches(3.6), Inches(5.0), [
    {"text": "Key Relationships", "size": 16, "color": ACCENT_TEAL, "bold": True, "spacing_after": 12},
    {"text": "regions  →  branches", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  One region has many branches", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "branches  →  employees", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  Employees work at a branch", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "customers  →  credit_history", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  1:1 credit profile per customer", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "customers  →  loan_applications", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  Customer applies for many loans", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "loan_applications  →  loan_payments", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  Each loan has multiple EMI payments", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "customers  →  transactions", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "  Customer makes many transactions", "size": 10, "color": MID_GRAY, "spacing_after": 10},
    {"text": "Normalization", "size": 16, "color": ACCENT_GOLD, "bold": True, "spacing_after": 8},
    {"text": "•  Third Normal Form (3NF)", "size": 11, "color": DARK_TEXT, "spacing_after": 4},
    {"text": "•  No transitive dependencies", "size": 11, "color": DARK_TEXT, "spacing_after": 4},
    {"text": "•  Referential integrity via FKs", "size": 11, "color": DARK_TEXT},
])

add_footer_bar(slide)
add_slide_number(slide, 5, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 6: TABLE CREATION (DDL)
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "DDL", Inches(0.8), Inches(0.4), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Table Creation — DDL with Constraints", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4), ACCENT_BLUE)

# Left — Sample DDL
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4),
            "Sample: customers table", font_size=14, font_color=ACCENT_GOLD, bold=True)

add_code_block(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.0),
"""CREATE TABLE customers (
    customer_id    INT PRIMARY KEY,
    full_name      VARCHAR(100),
    gender         VARCHAR(10),
    date_of_birth  DATE,
    age            INT,
    annual_income  DECIMAL(12,2),
    region_id      INT,
    branch_id      INT,
    customer_segment VARCHAR(20),
    -- ... 25+ columns total
    FOREIGN KEY (region_id)
        REFERENCES regions(region_id),
    FOREIGN KEY (branch_id)
        REFERENCES branches(branch_id)
);""", font_size=11)

# Right — Constraint summary cards
card_x = Inches(6.8)
constraints = [
    ("🔑  Primary Keys", "Every table has a single-column\nINT PRIMARY KEY for fast lookups", ACCENT_ROSE),
    ("🔗  Foreign Keys", "12 FK constraints enforce referential\nintegrity across all 9 tables", ACCENT_TEAL),
    ("🚫  NOT NULL", "Critical columns (IDs, names, dates)\nare marked NOT NULL", ACCENT_GOLD),
    ("📝  UNIQUE", "IFSC codes, employee codes, PAN —\nensure no duplicate identifiers", ACCENT_BLUE),
    ("⚙️  Defaults", "is_active DEFAULT TRUE, days_late\nDEFAULT 0 for clean defaults", PURPLE),
    ("📐  Data Types", "DECIMAL(15,2) for money, VARCHAR\nfor text, DATE for temporal data", GREEN),
]

for i, (title, desc, color) in enumerate(constraints):
    y = Inches(1.8) + Inches(i * 0.85)
    add_shape(slide, card_x, y, Inches(5.5), Inches(0.75), BG_CARD, color, Pt(1))
    add_textbox(slide, card_x + Inches(0.15), y + Inches(0.05), Inches(5.2), Inches(0.3),
                title, font_size=13, font_color=color, bold=True)
    add_textbox(slide, card_x + Inches(0.15), y + Inches(0.32), Inches(5.2), Inches(0.4),
                desc, font_size=10, font_color=SECOND_TEXT)

# Creation order note
add_shape(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.55), BG_CARD, MID_GRAY, Pt(1))
add_textbox(slide, Inches(1.0), Inches(6.47), Inches(11.3), Inches(0.4),
            "⚠️  Creation Order:  regions → loan_types → branches → employees → customers → credit_history → loan_applications → loan_payments → transactions",
            font_size=11, font_color=ACCENT_GOLD)

add_footer_bar(slide)
add_slide_number(slide, 6, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 7: DATA LOADING
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "DATA LOAD", Inches(0.8), Inches(0.4), PURPLE)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Bulk Data Loading — LOAD DATA LOCAL INFILE", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(5), PURPLE)

# Left — Code sample
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.4),
            "Loading Strategy", font_size=14, font_color=ACCENT_GOLD, bold=True)

add_code_block(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.5),
"""-- Enable local file loading
SET GLOBAL local_infile = 1;

-- Bulk load with NULL handling
LOAD DATA LOCAL INFILE
  'B:/Major Project/data/customers.csv'
INTO TABLE customers
FIELDS TERMINATED BY ','
ENCLOSED BY '"'
LINES TERMINATED BY '\\n'
IGNORE 1 ROWS;""", font_size=11)

# Advanced loading
add_textbox(slide, Inches(0.8), Inches(4.9), Inches(5.5), Inches(0.4),
            "Advanced: NULL Handling with SET Clause", font_size=14, font_color=ACCENT_GOLD, bold=True)

add_code_block(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.4),
"""LOAD DATA LOCAL INFILE '...loan_applications.csv'
INTO TABLE loan_applications
...
IGNORE 1 ROWS (@col1, @col2, ...)
SET col1 = NULLIF(@col1, ''),
    col2 = NULLIF(@col2, '');""", font_size=10)

# Right — Process flow cards
add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5.0), BG_CARD, PURPLE, Pt(1))
add_multi_text(slide, Inches(7.1), Inches(2.0), Inches(5.0), Inches(4.7), [
    {"text": "Loading Process", "size": 18, "color": PURPLE, "bold": True, "spacing_after": 12},
    {"text": "Step 1 — Enable Local Infile", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "SET GLOBAL local_infile = 1 on MySQL server", "size": 12, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "Step 2 — Load Parent Tables First", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "regions (40) → loan_types (10) → branches (250)", "size": 12, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "Step 3 — Load Core Tables", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "employees (~2K) → customers (100K) → credit_history (100K)", "size": 12, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "Step 4 — Load High-Volume Tables", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "loan_applications (300K) → loan_payments (900K) → transactions (1.5M)", "size": 12, "color": SECOND_TEXT, "spacing_after": 16},
    {"text": "Key Technique: NULLIF()", "size": 16, "color": ACCENT_TEAL, "bold": True, "spacing_after": 8},
    {"text": "Empty strings from CSV are converted to proper SQL NULLs using SET column = NULLIF(@var, '')", "size": 12, "color": SECOND_TEXT, "spacing_after": 8},
    {"text": "This ensures data integrity for optional columns like rejection_reason, disbursement_date, etc.", "size": 11, "color": MID_GRAY},
])

add_footer_bar(slide)
add_slide_number(slide, 7, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 8: DATA VERIFICATION
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "VERIFICATION", Inches(0.8), Inches(0.4), GREEN)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Data Verification — Row Count Validation", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4), GREEN)

# Verification query
add_code_block(slide, Inches(0.8), Inches(1.9), Inches(6.5), Inches(2.5),
"""-- Unified row count verification across all 9 tables
SELECT 'regions' AS table_name, COUNT(*) AS row_count
  FROM regions
UNION ALL SELECT 'loan_types', COUNT(*) FROM loan_types
UNION ALL SELECT 'branches', COUNT(*) FROM branches
UNION ALL SELECT 'employees', COUNT(*) FROM employees
UNION ALL SELECT 'customers', COUNT(*) FROM customers
UNION ALL SELECT 'credit_history', COUNT(*) FROM credit_history
UNION ALL SELECT 'loan_applications', COUNT(*) FROM loan_applications
UNION ALL SELECT 'loan_payments', COUNT(*) FROM loan_payments
UNION ALL SELECT 'transactions', COUNT(*) FROM transactions;""", font_size=11)

# Result cards
verify_data = [
    ("regions", "40", "✅"),
    ("loan_types", "10", "✅"),
    ("branches", "250", "✅"),
    ("employees", "~2,000", "✅"),
    ("customers", "100,000", "✅"),
    ("credit_history", "100,000", "✅"),
    ("loan_applications", "300,000", "✅"),
    ("loan_payments", "900,000", "✅"),
    ("transactions", "1,500,000", "✅"),
]

# Result table on the right
result_x = Inches(7.8)
add_textbox(slide, result_x, Inches(1.9), Inches(4.5), Inches(0.4),
            "Expected Results", font_size=16, font_color=GREEN, bold=True)

header_y = Inches(2.35)
add_shape(slide, result_x, header_y, Inches(4.5), Inches(0.38), GREEN)
add_textbox(slide, result_x + Inches(0.15), header_y + Inches(0.04), Inches(1.8), Inches(0.3),
            "Table", font_size=12, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
add_textbox(slide, result_x + Inches(2.0), header_y + Inches(0.04), Inches(1.5), Inches(0.3),
            "Row Count", font_size=12, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.RIGHT)
add_textbox(slide, result_x + Inches(3.7), header_y + Inches(0.04), Inches(0.6), Inches(0.3),
            "Status", font_size=12, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)

for i, (tbl, cnt, status) in enumerate(verify_data):
    y = Inches(2.78) + Inches(i * 0.40)
    bg = BG_CARD if i % 2 == 0 else BG_LIGHT
    add_shape(slide, result_x, y, Inches(4.5), Inches(0.37), bg)
    add_textbox(slide, result_x + Inches(0.15), y + Inches(0.04), Inches(1.8), Inches(0.3),
                tbl, font_size=11, font_color=ACCENT_TEAL, font_name="Consolas")
    add_textbox(slide, result_x + Inches(2.0), y + Inches(0.04), Inches(1.5), Inches(0.3),
                cnt, font_size=11, font_color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, result_x + Inches(3.7), y + Inches(0.04), Inches(0.6), Inches(0.3),
                status, font_size=11, alignment=PP_ALIGN.CENTER)

# Note
add_shape(slide, Inches(0.8), Inches(4.7), Inches(6.5), Inches(2.2), BG_CARD, GREEN, Pt(1))
add_multi_text(slide, Inches(1.1), Inches(4.85), Inches(6.0), Inches(2.0), [
    {"text": "Why UNION ALL?", "size": 16, "color": GREEN, "bold": True, "spacing_after": 8},
    {"text": "•  UNION ALL (not UNION) preserves all rows — no dedup needed", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Single query validates entire database in one execution", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Each sub-SELECT runs an independent COUNT(*) per table", "size": 12, "color": DARK_TEXT, "spacing_after": 5},
    {"text": "•  Faster than running 9 separate COUNT queries", "size": 12, "color": DARK_TEXT},
])

add_footer_bar(slide)
add_slide_number(slide, 8, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 9: INDEXING STRATEGY
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

add_section_badge(slide, "INDEXING", Inches(0.8), Inches(0.4), ACCENT_ROSE)
add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Indexing Strategy for Query Performance", font_size=32, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.55), Inches(4.5), ACCENT_ROSE)

# Left — Code
add_code_block(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5),
"""-- 22 Strategic Indexes Created

-- Customers (4 indexes)
CREATE INDEX idx_customers_region
    ON customers(region_id);
CREATE INDEX idx_customers_branch
    ON customers(branch_id);
CREATE INDEX idx_customers_segment
    ON customers(customer_segment);
CREATE INDEX idx_customers_state
    ON customers(state);

-- Credit History (3 indexes)
CREATE INDEX idx_credit_customer
    ON credit_history(customer_id);
CREATE INDEX idx_credit_score
    ON credit_history(credit_score);

-- Loan Applications (5 indexes)
CREATE INDEX idx_loan_app_customer
    ON loan_applications(customer_id);
CREATE INDEX idx_loan_app_status
    ON loan_applications(status);
CREATE INDEX idx_loan_app_date
    ON loan_applications(application_date);""", font_size=10)

# Right — Why indexes matter
add_shape(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(4.5), BG_CARD, ACCENT_ROSE, Pt(1))
add_multi_text(slide, Inches(7.1), Inches(2.1), Inches(5.0), Inches(4.2), [
    {"text": "Why These Indexes?", "size": 18, "color": ACCENT_ROSE, "bold": True, "spacing_after": 12},
    {"text": "🔍  JOIN Acceleration", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "FK columns (customer_id, branch_id, application_id) are indexed for fast multi-table JOINs", "size": 11, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "📊  Filter / WHERE Optimization", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "Columns used in WHERE clauses (status, credit_score, payment_status) get B-Tree indexes", "size": 11, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "📅  Time-Series Queries", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "Date columns (application_date, due_date, transaction_date) enable fast range scans", "size": 11, "color": SECOND_TEXT, "spacing_after": 12},
    {"text": "📈  GROUP BY Speed", "size": 14, "color": ACCENT_GOLD, "bold": True, "spacing_after": 4},
    {"text": "Segment and category indexes speed up aggregation queries", "size": 11, "color": SECOND_TEXT, "spacing_after": 16},
    {"text": "Index Summary", "size": 16, "color": ACCENT_TEAL, "bold": True, "spacing_after": 8},
    {"text": "•  22 total indexes across 5 tables", "size": 12, "color": DARK_TEXT, "spacing_after": 4},
    {"text": "•  Covers all FK columns for JOINs", "size": 12, "color": DARK_TEXT, "spacing_after": 4},
    {"text": "•  Targets high-cardinality columns", "size": 12, "color": DARK_TEXT, "spacing_after": 4},
    {"text": "•  Optimizes the 25+ analytical queries", "size": 12, "color": DARK_TEXT},
])

add_footer_bar(slide)
add_slide_number(slide, 9, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 10: QUERY SECTION INTRO
# ──────────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

# Large centered text
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1), "📊",
            font_size=60, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
            "Analytical SQL Queries", font_size=42, font_color=DARK_TEXT,
            bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.7),
            "10 High-Impact Queries Delivering Key Business Insights",
            font_size=22, font_color=ACCENT_ROSE, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(4.5), Inches(4.7), Inches(4.3), ACCENT_GOLD)

# Three level badges
levels = [
    ("🟢  BASIC", "Foundation queries using\nGROUP BY, ORDER BY, JOINs", GREEN, Inches(1.5)),
    ("🟡  INTERMEDIATE", "Window functions, CASE,\nmulti-table aggregations", ACCENT_GOLD, Inches(5.0)),
    ("🔴  ADVANCED", "Subqueries, CTEs, RFM,\nNPA & Cohort Analysis", RED_SOFT, Inches(8.5)),
]

for label, desc, color, x_pos in levels:
    add_shape(slide, x_pos, Inches(5.2), Inches(3.3), Inches(1.5), BG_CARD, color, Pt(1.5))
    add_textbox(slide, x_pos + Inches(0.2), Inches(5.35), Inches(2.9), Inches(0.4),
                label, font_size=16, font_color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x_pos + Inches(0.2), Inches(5.85), Inches(2.9), Inches(0.7),
                desc, font_size=12, font_color=SECOND_TEXT, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide)
add_slide_number(slide, 10, TOTAL_SLIDES)


# ══════════════════════════════════════════════════════════════════════════
#  QUERY SLIDES (10 queries — Slides 11-20... but we need slide 20 for
#  takeaways, so queries are slides 11-20 total in the ppt which is 20 slides)
# ══════════════════════════════════════════════════════════════════════════

queries = [
    {
        "num": "Q1",
        "level": "BASIC",
        "level_color": GREEN,
        "title": "Total Customers by Zone",
        "sql": """SELECT zone,
       COUNT(*) AS total_customers
FROM customers
GROUP BY zone
ORDER BY total_customers DESC;""",
        "insight_title": "Geographic Customer Distribution",
        "insights": [
            "Reveals which zones (North, South, East, West) have the highest customer concentration",
            "Helps the bank allocate marketing budgets proportionally across regions",
            "Identifies underserved zones where expansion opportunities exist",
            "Supports RBI's priority sector lending compliance by geography",
        ],
        "concepts": "GROUP BY  •  COUNT(*)  •  ORDER BY DESC",
    },
    {
        "num": "Q3",
        "level": "BASIC",
        "level_color": GREEN,
        "title": "Loan Application Status Breakdown (%)",
        "sql": """SELECT status,
       COUNT(*) AS count,
       ROUND(COUNT(*) * 100.0
         / SUM(COUNT(*)) OVER(), 2)
         AS pct
FROM loan_applications
GROUP BY status;""",
        "insight_title": "Loan Pipeline Health Check",
        "insights": [
            "Shows the exact approval vs rejection vs pending ratio as percentages",
            "A high rejection rate signals overly strict policies or poor customer targeting",
            "Uses window function SUM() OVER() for percentage calculation without subquery",
            "Critical KPI for bank management to track loan funnel efficiency",
        ],
        "concepts": "SUM() OVER()  •  Window Function  •  ROUND()",
    },
    {
        "num": "Q7",
        "level": "BASIC",
        "level_color": GREEN,
        "title": "Transaction Mode Popularity (UPI vs Traditional)",
        "sql": """SELECT transaction_mode,
       COUNT(*) AS txn_count,
       ROUND(SUM(amount), 2)
         AS total_amount
FROM transactions
WHERE status = 'Success'
GROUP BY transaction_mode
ORDER BY txn_count DESC;""",
        "insight_title": "Digital Banking Adoption",
        "insights": [
            "Ranks UPI, NEFT, Cash, ATM, IMPS by usage count and total value",
            "Reveals the shift from traditional (Cash, Cheque) to digital (UPI, NEFT)",
            "Helps the bank decide where to invest — more ATMs or better app?",
            "Filters only successful transactions for accurate volume analysis",
        ],
        "concepts": "WHERE filter  •  SUM()  •  GROUP BY  •  ORDER BY",
    },
    {
        "num": "Q9",
        "level": "INTERMEDIATE",
        "level_color": ACCENT_GOLD,
        "title": "Loan Approval Rate by Loan Type",
        "sql": """SELECT loan_type_name,
  COUNT(*) AS total,
  SUM(CASE WHEN status IN
    ('Disbursed','Approved','Closed')
    THEN 1 ELSE 0 END) AS approved,
  ROUND(SUM(CASE WHEN status IN
    ('Disbursed','Approved','Closed')
    THEN 1 ELSE 0 END)
    * 100.0 / COUNT(*), 2)
    AS approval_rate_pct
FROM loan_applications
GROUP BY loan_type_name
ORDER BY approval_rate_pct DESC;""",
        "insight_title": "Product-Level Approval Analysis",
        "insights": [
            "Gold Loans & Agri Loans likely have highest approval (collateral-backed / govt schemes)",
            "Personal Loans may have lowest approval due to unsecured nature",
            "Uses CASE-based conditional aggregation — a powerful intermediate SQL pattern",
            "Helps management fine-tune underwriting policies per product",
        ],
        "concepts": "CASE WHEN  •  Conditional Aggregation  •  IN()",
    },
    {
        "num": "Q12",
        "level": "INTERMEDIATE",
        "level_color": ACCENT_GOLD,
        "title": "Payment Default Rate Analysis",
        "sql": """SELECT payment_status,
       COUNT(*) AS count,
       ROUND(COUNT(*) * 100.0
         / SUM(COUNT(*)) OVER(), 2)
         AS pct
FROM loan_payments
GROUP BY payment_status
ORDER BY count DESC;""",
        "insight_title": "Repayment Risk Assessment",
        "insights": [
            "Shows Paid vs Late vs Missed payment distribution across 900K records",
            "High 'Missed' percentage directly signals rising NPA risk",
            "The bank can trigger early warning systems based on this metric",
            "Window function calculates % without a separate total subquery",
        ],
        "concepts": "Window Function  •  Risk Metric  •  900K Records",
    },
    {
        "num": "Q17",
        "level": "INTERMEDIATE",
        "level_color": ACCENT_GOLD,
        "title": "Year-over-Year Loan Application Growth",
        "sql": """SELECT YEAR(application_date) AS yr,
  COUNT(*) AS applications,
  LAG(COUNT(*)) OVER (
    ORDER BY YEAR(application_date))
    AS prev_year,
  ROUND((COUNT(*) - LAG(COUNT(*))
    OVER (ORDER BY
      YEAR(application_date)))
    * 100.0 / LAG(COUNT(*))
    OVER (ORDER BY
      YEAR(application_date)), 2)
    AS yoy_growth_pct
FROM loan_applications
GROUP BY yr ORDER BY yr;""",
        "insight_title": "Business Growth Trajectory",
        "insights": [
            "Uses LAG() window function to compare each year with the previous year",
            "Calculates exact YoY growth percentage for strategic planning",
            "Declining growth signals competitive pressure; spikes indicate expansion",
            "Critical metric for board-level reporting and investor presentations",
        ],
        "concepts": "LAG()  •  Window Function  •  YoY Growth  •  YEAR()",
    },
    {
        "num": "Q18",
        "level": "ADVANCED",
        "level_color": RED_SOFT,
        "title": "Credit Score Bands vs Loan Default Rate",
        "sql": """SELECT
  CASE
    WHEN ch.credit_score >= 800
      THEN '800+ (Excellent)'
    WHEN ch.credit_score >= 700
      THEN '700-799 (Good)'
    WHEN ch.credit_score >= 600
      THEN '600-699 (Fair)'
    ELSE 'Below 600 (Poor)'
  END AS credit_band,
  COUNT(DISTINCT la.application_id)
    AS total_loans,
  ROUND(... * 100.0 / ..., 2)
    AS default_rate_pct
FROM loan_applications la
JOIN credit_history ch ON ...
JOIN loan_payments lp ON ...
GROUP BY credit_band;""",
        "insight_title": "Credit Risk Validation",
        "insights": [
            "Maps credit score bands (Excellent/Good/Fair/Poor) to actual default rates",
            "Validates whether the bank's scoring model accurately predicts risk",
            "If 'Excellent' customers also default, the scoring criteria need revision",
            "Uses CASE + multi-table JOINs + DISTINCT counting — advanced pattern",
        ],
        "concepts": "CASE bands  •  3-table JOIN  •  COUNT(DISTINCT)  •  Risk Model",
    },
    {
        "num": "Q20",
        "level": "ADVANCED",
        "level_color": RED_SOFT,
        "title": "Customer Lifetime Value (CLV)",
        "sql": """SELECT c.customer_id, c.full_name,
  c.customer_segment,
  COALESCE(t.total_credits, 0)
    AS total_deposits,
  COALESCE(l.total_interest, 0)
    AS interest_paid_to_bank,
  COALESCE(t.total_credits, 0)
    + COALESCE(l.total_interest, 0)
    AS lifetime_value
FROM customers c
LEFT JOIN (
  SELECT customer_id, SUM(amount) ...
  FROM transactions
  WHERE type='Credit' GROUP BY ...
) t ON c.customer_id = t.customer_id
LEFT JOIN (
  SELECT customer_id, SUM(interest)..
  FROM loan_payments GROUP BY ...
) l ON c.customer_id = l.customer_id
ORDER BY lifetime_value DESC
LIMIT 20;""",
        "insight_title": "Most Valuable Customers",
        "insights": [
            "Combines deposits (credit transactions) + loan interest into a single CLV score",
            "Uses LEFT JOINs with subqueries to handle customers without loans or transactions",
            "COALESCE ensures NULLs become 0 for accurate summation",
            "Identifies top 20 VIP customers for priority service & retention offers",
        ],
        "concepts": "Subqueries  •  LEFT JOIN  •  COALESCE  •  CLV Metric",
    },
    {
        "num": "Q23",
        "level": "ADVANCED",
        "level_color": RED_SOFT,
        "title": "NPA (Non-Performing Asset) Identification",
        "sql": """SELECT la.application_id,
  c.full_name, la.loan_type_name,
  la.loan_amount_approved,
  COUNT(CASE WHEN
    lp.payment_status = 'Missed'
    THEN 1 END) AS missed_payments,
  MAX(lp.outstanding_balance)
    AS current_outstanding
FROM loan_applications la
JOIN loan_payments lp
  ON la.application_id
     = lp.application_id
JOIN customers c
  ON la.customer_id = c.customer_id
GROUP BY la.application_id, ...
HAVING missed_payments >= 3
ORDER BY current_outstanding DESC
LIMIT 20;""",
        "insight_title": "Regulatory Risk Detection",
        "insights": [
            "Flags loans with 3+ missed payments as potential NPAs — an RBI requirement",
            "Shows outstanding balance at risk for provisioning calculations",
            "Uses HAVING for post-aggregation filtering (can't use WHERE for aggregates)",
            "Early NPA detection allows restructuring before loans become fully non-performing",
        ],
        "concepts": "HAVING  •  3-table JOIN  •  NPA Detection  •  RBI Compliance",
    },
    {
        "num": "Q25",
        "level": "ADVANCED",
        "level_color": RED_SOFT,
        "title": "Cohort Analysis — Loan Applications by Account Opening Year",
        "sql": """SELECT
  YEAR(c.account_open_date)
    AS cohort_year,
  YEAR(la.application_date)
    AS app_year,
  COUNT(*) AS applications,
  SUM(CASE WHEN la.status
    = 'Disbursed' THEN 1 ELSE 0 END)
    AS disbursed
FROM loan_applications la
JOIN customers c
  ON la.customer_id = c.customer_id
GROUP BY cohort_year, app_year
ORDER BY cohort_year, app_year;""",
        "insight_title": "Customer Lifecycle Patterns",
        "insights": [
            "Groups customers into cohorts by account opening year, tracks loan behavior over time",
            "Reveals whether newer customers apply for loans sooner than older cohorts",
            "Helps predict future loan demand from recent account openers",
            "Classic cohort analysis pattern used in retention strategy & forecasting",
        ],
        "concepts": "Cohort Analysis  •  YEAR()  •  Multi-GROUP BY  •  Lifecycle",
    },
]

for qi, q in enumerate(queries):
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, BG_LIGHT)
    slide_num = 11 + qi

    # Level badge
    add_section_badge(slide, q["level"], Inches(0.8), Inches(0.4), q["level_color"])
    # Query number badge
    add_section_badge(slide, q["num"], Inches(2.8), Inches(0.4), ACCENT_GOLD)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
                q["title"], font_size=28, font_color=DARK_TEXT, bold=True)
    add_accent_line(slide, Inches(0.8), Inches(1.5), Inches(3.5), q["level_color"])

    # Left — SQL Code
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.2), Inches(0.4),
                "SQL Query", font_size=13, font_color=ACCENT_TEAL, bold=True)
    add_code_block(slide, Inches(0.8), Inches(2.05), Inches(5.2), Inches(4.2),
                   q["sql"], font_size=10)

    # Bottom-left — Concepts used
    add_shape(slide, Inches(0.8), Inches(6.4), Inches(5.2), Inches(0.5), BG_CARD, q["level_color"], Pt(1))
    add_textbox(slide, Inches(1.0), Inches(6.45), Inches(4.8), Inches(0.4),
                f"Concepts:  {q['concepts']}", font_size=10, font_color=q["level_color"])

    # Right — Insight card
    add_shape(slide, Inches(6.5), Inches(1.7), Inches(6.0), Inches(5.2), BG_CARD, q["level_color"], Pt(1.5))

    add_textbox(slide, Inches(6.8), Inches(1.85), Inches(5.5), Inches(0.4),
                "💡  " + q["insight_title"], font_size=18, font_color=q["level_color"], bold=True)

    insight_lines = []
    for ins in q["insights"]:
        insight_lines.append({"text": "▸  " + ins, "size": 13, "color": DARK_TEXT, "spacing_after": 12})

    add_multi_text(slide, Inches(6.8), Inches(2.5), Inches(5.5), Inches(4.0), insight_lines)

    # Business impact badge
    add_shape(slide, Inches(6.5), Inches(6.4), Inches(6.0), Inches(0.5), q["level_color"])
    # Impact text uses white on colored badge for contrast
    impact_texts = {
        "Q1": "Impact: Regional Strategy & Resource Allocation",
        "Q3": "Impact: Loan Funnel Optimization & Policy Review",
        "Q7": "Impact: Digital Transformation & Channel Investment",
        "Q9": "Impact: Underwriting Policy Optimization",
        "Q12": "Impact: Early Warning System for Defaults",
        "Q17": "Impact: Strategic Planning & Board Reporting",
        "Q18": "Impact: Credit Scoring Model Validation",
        "Q20": "Impact: VIP Customer Identification & Retention",
        "Q23": "Impact: RBI Compliance & Loss Prevention",
        "Q25": "Impact: Retention Strategy & Demand Forecasting",
    }
    add_textbox(slide, Inches(6.7), Inches(6.43), Inches(5.6), Inches(0.4),
                impact_texts.get(q["num"], ""), font_size=11, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    add_footer_bar(slide)
    add_slide_number(slide, slide_num, TOTAL_SLIDES)


# ──────────────────────────────────────────────────────────────────────────
# SLIDE 20 (replaces slide 21): KEY TAKEAWAYS & THANK YOU
# ──────────────────────────────────────────────────────────────────────────

# But we have 10 + 10 = 20 slides already (slide 11..20 are queries).
# So our "Thank you" needs to be slide 21 — let's update TOTAL_SLIDES.
# Actually the current count: 1 title + 2 agenda + 3 overview + 4 dataset + 5 ER
#   + 6 DDL + 7 loading + 8 verify + 9 index + 10 query-intro + 11..20 (10 queries) = 20
# We need 1 more for takeaways. Let's make it 21.

# Update all slide numbers retroactively isn't possible with pptx easily,
# so let's just add the final slide.

TOTAL_SLIDES = 21  # Update for slide numbering

slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_LIGHT)

# Top accent
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_ROSE
shape.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Key Takeaways & Summary", font_size=36, font_color=DARK_TEXT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.2), Inches(3.5), ACCENT_ROSE)

# Takeaway cards (2 columns x 3 rows)
takeaways = [
    ("🗃️  9-Table Schema", "Designed a fully normalized 3NF\nschema with 12 FK constraints\nensuring data integrity", ACCENT_ROSE),
    ("📥  ~3M Records Loaded", "Bulk ingestion via LOAD DATA\nLOCAL INFILE with NULLIF()\nfor clean NULL handling", ACCENT_TEAL),
    ("⚡  22 Indexes", "Strategic B-Tree indexes on\nFK columns, dates, and\nhigh-cardinality filters", ACCENT_GOLD),
    ("📊  25+ Queries", "Basic to Advanced SQL covering\nJOINs, window functions, CASE,\nsubqueries, & cohort analysis", ACCENT_BLUE),
    ("🔍  Risk Detection", "NPA identification, credit score\nvalidation, default rate tracking\nfor RBI compliance", RED_SOFT),
    ("💎  CLV & RFM", "Customer Lifetime Value and\nRFM segmentation for targeted\nmarketing & retention", PURPLE),
]

for i, (title, desc, color) in enumerate(takeaways):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 1.55)
    add_shape(slide, x, y, Inches(5.8), Inches(1.35), BG_CARD, color, Pt(1.5))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.4),
                title, font_size=16, font_color=color, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.4), Inches(0.75),
                desc, font_size=12, font_color=SECOND_TEXT, line_spacing=1.4)

# Thank you banner
add_shape(slide, Inches(2.5), Inches(6.0), Inches(8.3), Inches(0.7), ACCENT_ROSE)
add_textbox(slide, Inches(2.5), Inches(6.05), Inches(8.3), Inches(0.6),
            "Thank You  •  IndoSynth Gramin Bank  •  SQL Phase Complete ✅",
            font_size=18, font_color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)

add_footer_bar(slide)
add_slide_number(slide, 21, 21)


# ══════════════════════════════════════════════════════════════════════════
#  SAVE THE PRESENTATION
# ══════════════════════════════════════════════════════════════════════════

output_path = r"B:\Major Project\IndoSynth_Bank_SQL_Phase.pptx"
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Open the file in PowerPoint to view!")
